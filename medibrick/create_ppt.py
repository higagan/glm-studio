from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(20)
        p.space_after = Pt(12)
    
    return slide

# Slide 1: Title
add_title_slide(prs, "Medibrick Cofounder Meeting", "Gagan (Tech) + Arpana (Medical)\nJune 2026")

# Slide 2: Meeting Goals
add_content_slide(prs, "Meeting Goals", [
    "• Align on current platform status",
    "• Decide MVP scope for first hospital pilot",
    "• Define roles: Gagan vs Arpana",
    "• Set 2-week action plan with owners",
    "",
    "Outcome: Clear decisions + immediate next steps"
])

# Slide 3: Current Status
add_content_slide(prs, "Platform Status: What's Built", [
    "✅ Landing page (medibrick.com)",
    "✅ Database (Supabase PostgreSQL)",
    "🟡 Basic auth (partial)",
    "🔴 Shift posting (not started)",
    "🔴 Application flow (not started)",
    "🔴 Payment escrow (not started)",
    "🔴 Review system (not started)",
    "",
    "Decision: What's minimum to test with 1 hospital?"
])

# Slide 4: Pain Point 1 - No-Show
add_content_slide(prs, 'Pain Point #1: "Yes but No-Show"', [
    "Problem: Nurse confirms shift → Doesn't show → Hospital short-staffed",
    "",
    "Why it happens:",
    "• Better offer elsewhere",
    "• Multiple bookings (picked best)",
    "• Zero penalty for no-show",
    "",
    "Medibrick Solution:",
    "✅ Verified profiles + deposit (₹500-2000)",
    "✅ Calendar lock (prevents double-booking)",
    "✅ Instant backup matching (#2 auto-promoted)",
    "✅ 3 no-shows = profile suspended",
    "",
    "Target: 95% show rate (industry: ~70%)"
])

# Slide 5: Pain Point 2 - Payment
add_content_slide(prs, "Pain Point #2: Delayed Payments", [
    "Problem: Professional completes shift → Waits 2-4 weeks → Chases admin",
    "",
    "Why it happens:",
    "• Hospital pays platform late",
    "• Platform holds money as working capital",
    "• Manual invoicing, no transparency",
    "",
    "Medibrick Solution:",
    "✅ Escrow: Hospital deposits payment upfront",
    "✅ Auto-release: 24h after shift completion",
    "✅ Dispute window: 48h for hospital",
    "✅ Live tracking: Pending → Processing → Paid",
    "",
    "Result: Happy professionals, repeat usage"
])

# Slide 6: Jobizo vs Medibrick
add_content_slide(prs, "Competitor: Jobizo.com", [
    "Jobizo: 10M+ shifts, 150+ hospitals, 60K+ professionals",
    "",
    "Their Weaknesses = Our Opportunities:",
    "",
    "• Job-seeker focused → We are institution-first",
    "• Allopathic doctors only → We add AYUSH + nurses + techs",
    "• No Ayurvedic coverage → 700K+ practitioners, zero competition",
    "• Hospitals only → We serve clinics + diagnostic + wellness",
    "• No no-show protection → Verified + deposit + backup = 95%",
    "• Payment: 2-4 weeks → We pay in 24 hours",
    "",
    "We don't compete head-to-head. We own the gaps."
])

# Slide 7: Target Market
add_content_slide(prs, "Target Market: Phase 1", [
    "Cities (in order):",
    "1. Bengaluru (HQ) - Start here",
    "2. Hyderabad",
    "3. Chennai",
    "4. Mumbai",
    "5. Delhi",
    "6. Pune",
    "",
    "Segments Jobizo ignores:",
    "• Ayurvedic centers (700K+ practitioners)",
    "• Diagnostic chains (Thyrocare, SRL, Metropolis)",
    "• Nursing homes (not big hospitals)",
    "• Small clinics (10-50 beds)"
])

# Slide 8: Business Model
add_content_slide(prs, "Business Model", [
    "Recommended: Hybrid Pricing",
    "",
    "Free tier:",
    "• Free to post shifts",
    "• ₹200 per filled shift",
    "",
    "Premium tier (₹5000/month):",
    "• Unlimited shifts",
    "• Priority matching",
    "• Analytics dashboard",
    "",
    "Professional side:",
    "• Free to join, free to apply",
    "• Deposit: ₹500-2000 (refundable)",
    "",
    "Pilot offer: First 3 shifts FREE"
])

# Slide 9: Arpana's Role
add_content_slide(prs, "Arpana's Medical Advantage", [
    "Why Arpana's background is critical:",
    "",
    "• Hospital Trust: Doctors trust doctors",
    "  → Opens doors that cold outreach can't",
    "",
    "• Medical Accuracy: Speaks the right language",
    "  → Departments, certifications, shift types",
    "",
    "• Professional Recruitment: Her network",
    "  → Seed 50+ profiles before launch",
    "",
    "• Clinical Validation: What actually matters",
    "  → vs what's nice-to-have",
    "",
    "• Regulatory Guidance: Licensing, compliance",
    "  → Navigate medical council requirements"
])

# Slide 10: Next 2 Weeks
add_content_slide(prs, "Action Plan: Next 2 Weeks", [
    "Week 1:",
    "• Mon: Secure database (RLS, headers)",
    "• Tue: Decide MVP scope (Gagan + Arpana)",
    "• Wed: Create 50 seed professional profiles (Arpana)",
    "• Thu: Outreach to 5 hospitals (Arpana leads)",
    "• Fri: Apply to IIHMR + NSRCEL accelerators (Gagan)",
    "",
    "Week 2:",
    "• Hospital meetings + feedback (Both)",
    "• Set up Razorpay escrow prototype (Gagan)",
    "• Revise pitch based on feedback (Both)",
    "• WhatsApp notification system (Gagan)"
])

# Slide 11: Decisions Needed
add_content_slide(prs, "Decisions Needed Today", [
    "1. MVP Scope:",
    "   ☐ Full platform (4-6 weeks)",
    "   ☐ Basic + WhatsApp alerts (1-2 weeks)",
    "   ☐ Pure manual (this week)",
    "",
    "2. First Segment:",
    "   ☐ Large hospitals (slow, high volume)",
    "   ☐ Small clinics (fast, lower volume)",
    "   ☐ Ayurvedic centers (no competition)",
    "",
    "3. Roles:",
    "   ☐ Gagan: Tech + Ops + Fundraising",
    "   ☐ Arpana: Medical + Sales + Network",
    "",
    "4. Timeline: First paid shift?",
    "   ☐ This month (aggressive)"
])

# Slide 12: Closing
add_title_slide(prs, "Let's Build This", "Questions? Decisions? Next steps?")

# Save
output_path = "/Users/gagandeep/.openclaw/workspace/medibrick/medibrick-cofounder-meeting.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
