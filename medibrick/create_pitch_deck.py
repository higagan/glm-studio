from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(20)
        p.space_after = Pt(12)
    
    return slide

# Slide 1: Title
add_title_slide(prs, "Medibrick", "Verified Marketplace for Healthcare Shifts\nGagan + Arpana | June 2026")

# Slide 2: Problem
add_content_slide(prs, "The Problem: Empty Shifts, Stressed Hospitals", [
    "• Hospitals face 30-40% no-show rates for temporary staff",
    "• Admin calls 10 people, hopes 1 shows up",
    "• No verification, no commitment, no backup",
    "",
    "• Professionals wait weeks for payment",
    "• Chasing finance, getting frustrated, leaving platforms",
    "",
    "Result: Hospitals are short-staffed. Professionals are unpaid. Patients wait."
])

# Slide 3: Solution
add_content_slide(prs, "Medibrick: Verified Locum Marketplace", [
    "For Hospitals:",
    "• Post shifts, get verified professionals who show up",
    "• Browse verified profiles with reviews from other hospitals",
    "• See professional's track record before selecting",
    "",
    "For Professionals:",
    "• Flexible shifts, fair pay, build reputation",
    "• Public profile with reviews helps get more shifts",
    "• Rate hospitals on payment speed and treatment",
    "",
    "Trust through transparency, not deposits."
])

# Slide 4: Market Opportunity
add_content_slide(prs, "Market Opportunity", [
    "India's Healthcare Staffing Gap:",
    "• 700,000+ Ayurvedic practitioners — NO platform serves them",
    "• 2.5M+ nurses — fragmented, offline hiring",
    "• Diagnostic chains need temp technicians for peak hours",
    "• Small clinics (10-50 beds) ignored by large platforms",
    "",
    "Phase 1 Cities: Bengaluru → Hyderabad → Chennai → Mumbai → Delhi → Pune",
    "",
    "TAM: ₹2,000 Cr+ (healthcare temp staffing in India)"
])

# Slide 5: Competitive Advantage
add_content_slide(prs, "Why Now? Why Us?", [
    "Jobizo (current leader):",
    "• Allopathic doctors only, job-seeker first, hospitals only",
    "• No review system, no transparency",
    "",
    "Medibrick:",
    "• Institution-first: Hospital posts, professionals apply",
    "• Allopathic + AYUSH + nurses + technicians",
    "• Hospitals + clinics + diagnostic + wellness + Ayurvedic",
    "• Public reviews build trust without forcing deposits",
    "• Hospitals handle payment directly — no platform delay",
    "",
    "We don't compete with Jobizo. We own the gaps they ignore."
])

# Slide 6: Business Model
add_content_slide(prs, "Business Model", [
    "Phase 1: Free to Onboard",
    "• Hospitals: Free to post shifts",
    "• Professionals: Free to join, free to apply",
    "• Goal: Build density of hospitals and professionals",
    "",
    "Phase 2: Monetize Trust",
    "• Premium listings for hospitals (featured shifts)",
    "• Verified badge for professionals (background check)",
    "• Analytics dashboard for hospital staffing patterns",
    "",
    "Payment: Hospitals pay professionals directly.",
    "Platform stays neutral — reviews keep everyone honest."
])

# Slide 7: Current Status
add_content_slide(prs, "Current Status", [
    "Built:",
    "• Landing page (medibrick.com)",
    "• Database (Supabase)",
    "• Shift posting + application flow",
    "• Public profile pages with reviews",
    "",
    "Pre-Launch:",
    "• Review system in development",
    "• Targeting first 5 Bengaluru hospitals this month",
    "• Onboarding first 50 verified professionals",
    "",
    "Team:",
    "• Gagan: Tech + Operations (Software Engineer)",
    "• Arpana: Medical Strategy + Hospital Sales (Healthcare Background)"
])

# Slide 8: What We're Looking For
add_content_slide(prs, "What We're Looking For", [
    "From This Accelerator:",
    "• Hospital Network: Introductions to Bengaluru hospitals",
    "• Mentorship: Healthcare industry experts",
    "• Credibility: Backing for hospital sales conversations",
    "• Learning: Best practices from other healthcare startups",
    "",
    "Use of Funds (if applicable):",
    "• Product development: Review system, mobile app",
    "• Hospital acquisition: Pilot program, onboarding support",
    "• Professional onboarding: Verify first 500 profiles",
    "• Operations: Team expansion"
])

# Slide 9: Vision
add_title_slide(prs, "Vision", "Every empty shift in India filled.\nEvery healthcare professional treated fairly.\n\nStarting with Bengaluru.")

# Save
output_path = "/Users/gagandeep/.openclaw/workspace/medibrick/Medibrick-Pitch-Deck.pptx"
prs.save(output_path)
print(f"Updated: {output_path}")
